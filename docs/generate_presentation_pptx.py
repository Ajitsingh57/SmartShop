import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    os.makedirs("d:/SmartShop/docs", exist_ok=True)
    pptx_path = "d:/SmartShop/docs/SmartShop_Project_Presentation.pptx"
    
    prs = Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette definition
    NAVY = RGBColor(15, 23, 42)        # Slate 900
    PRIMARY = RGBColor(30, 58, 138)     # Deep Blue
    ACCENT = RGBColor(2, 132, 199)      # Sky Blue 600
    LIGHT_BG = RGBColor(248, 250, 252)  # Slate 50
    CARD_BG = RGBColor(255, 255, 255)   # White
    CARD_BORDER = RGBColor(203, 213, 225) # Slate 300
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(22, 163, 74)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="MAJOR PROJECT PRESENTATION"):
        # Header banner
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY
        
        # Subtle horizontal divider
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.733), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = CARD_BORDER
        line.line.color.rgb = CARD_BORDER

    def add_card(slide, left, top, width, height, title, items, bg_color=CARD_BG, border_color=CARD_BORDER, title_color=PRIMARY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        
        tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = "Arial"
            p_title.font.size = Pt(15)
            p_title.font.bold = True
            p_title.font.color.rgb = title_color
            p_title.space_after = Pt(10)
        
        for i, item in enumerate(items):
            p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
            p.text = f"•  {item}" if not item.startswith(" ") else item
            p.font.name = "Arial"
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (Dark Premium Theme)
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.fill.background()

    # Title box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(3.0))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    
    p_tag = tf1.paragraphs[0]
    p_tag.text = "MAJOR PROJECT PRESENTATION"
    p_tag.font.name = "Arial"
    p_tag.font.size = Pt(14)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT
    
    p_main = tf1.add_paragraph()
    p_main.text = "SmartShop"
    p_main.font.name = "Arial"
    p_main.font.size = Pt(44)
    p_main.font.bold = True
    p_main.font.color.rgb = WHITE
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "A Modern Digital Shop Management, POS & Customer Khata Ledger Platform"
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = RGBColor(148, 163, 184)
    p_sub.space_before = Pt(8)

    # Info card on title slide
    meta_box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.7), Inches(11.333), Inches(1.8))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
    meta_box.line.color.rgb = RGBColor(51, 65, 85)
    
    mtb = s1.shapes.add_textbox(Inches(1.3), Inches(4.9), Inches(10.733), Inches(1.4))
    mtf = mtb.text_frame
    mtf.word_wrap = True
    
    p1 = mtf.paragraphs[0]
    p1.text = "Submitted By:  Ajit Singh & Team            |      Branch: Computer Science & Engineering"
    p1.font.name = "Arial"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = mtf.add_paragraph()
    p2.text = "Project Guide:  Faculty Guide / Project Coordinator"
    p2.font.name = "Arial"
    p2.font.size = Pt(13)
    p2.font.color.rgb = RGBColor(203, 213, 225)
    p2.space_before = Pt(8)
    
    p3 = mtf.add_paragraph()
    p3.text = "Academic Session: 2025 – 2026                 |      Tech Stack: MERN Stack (MongoDB, Express, React 19, Node.js)"
    p3.font.name = "Arial"
    p3.font.size = Pt(12)
    p3.font.color.rgb = ACCENT
    p3.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 2: PROJECT OVERVIEW & ABSTRACT
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Project Overview & Executive Summary")
    
    add_card(s2, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1), 
             "What is SmartShop?", [
                 "SmartShop is a cloud-based web application tailored for retail stores, supermarkets, and general merchants.",
                 "Combines high-speed Point-of-Sale (POS) counter billing with an integrated digital customer ledger ('Khata').",
                 "Eliminates paper registers by automatically recording customer dues, trust ratings, borrow limits, and settlements.",
                 "Includes full customer self-service access to monitor dues, view purchases, and submit payments."
             ])
             
    add_card(s2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1), 
             "Core Value Proposition", [
                 "Zero Data Loss: Replaces fragile paper bahi-khata with cloud storage and automated daily backups.",
                 "Transparent Customer Ledger: Real-time itemized bills visible directly on customers' mobile phones.",
                 "Automated Returns & Restocking: Full and partial item returns automatically adjust inventory and debts.",
                 "Multi-Channel Settlements: Supports Cash, manual UPI with screenshot verification, and Razorpay gateway."
             ])

    # ----------------------------------------------------
    # SLIDE 3: PROBLEM STATEMENT & MOTIVATION
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Problem Statement & Need for Digitalization")
    
    add_card(s3, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Traditional Khata Issues", [
                 "Manual arithmetic errors during daily ledger tallying.",
                 "Physical register damage, ink fading, or page loss.",
                 "Disputes over old unverified credit balances.",
                 "No automated notifications for overdue repayments."
             ], bg_color=RGBColor(254, 242, 242), border_color=RGBColor(252, 165, 165), title_color=RGBColor(185, 28, 28))
             
    add_card(s3, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Standalone POS Limitations", [
                 "Existing POS systems focus only on instant cash/card billing.",
                 "Cannot manage customer credit accounts or udhar history.",
                 "Require expensive proprietary hardware and setup.",
                 "Complex, non-intuitive interfaces for retail staff."
             ], bg_color=RGBColor(254, 243, 199), border_color=RGBColor(253, 230, 138), title_color=RGBColor(180, 83, 9))
             
    add_card(s3, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "The SmartShop Solution", [
                 "Unified POS + Digital Khata in a single web application.",
                 "Accessible on any laptop, tablet, or smartphone.",
                 "ACID-compliant transactional rollbacks for zero discrepancies.",
                 "Zero hardware lock-in and seamless cloud synchronization."
             ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208), title_color=GREEN)

    # ----------------------------------------------------
    # SLIDE 4: SYSTEM OBJECTIVES
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Project Objectives & Scope")
    
    add_card(s4, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Key Functional Objectives", [
                 "Fast Multi-Mode Billing: Enable counter staff to bill items in under 30 seconds via Cash, UPI, Credit, or Partial splits.",
                 "Credit Ceiling Enforcement: Validate customer borrow limits before issuing new credit to prevent bad debts.",
                 "Full & Partial Returns: Automatically restock items and deduct customer dues atomically upon product return.",
                 "Payment Claims Reconciliation: Provide admin verification dashboard for customer-submitted Cash & UPI proofs.",
                 "Audit Logging: Track admin actions (product updates, price changes, sales, and approvals) for accountability."
             ])
             
    add_card(s4, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Non-Functional & Technical Goals", [
                 "ACID Atomicity: Multi-document operations wrapped in MongoDB transactions to prevent dirty writes.",
                 "High Responsiveness: Modern glassmorphism UI optimized for Mobile, Tablet, and Desktop screens.",
                 "Robust Security: Stateless JWT authentication, role-based access control (RBAC), and security headers.",
                 "Zero-Config Media Storage: Cloudinary integration for scalable product and payment proof storage.",
                 "Scalable MERN Architecture: Clean Separation of Concerns with RESTful endpoints and modular React structure."
             ])

    # ----------------------------------------------------
    # SLIDE 5: TECHNOLOGY STACK
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Technology Stack & Tools")
    
    add_card(s5, Inches(0.8), Inches(1.7), Inches(2.7), Inches(5.1),
             "Frontend UI", [
                 "React 19 (Component UI)",
                 "Vite (Lightning Fast Build)",
                 "TailwindCSS (Styling)",
                 "Lucide React (Icons)",
                 "Glassmorphism Design"
             ])
             
    add_card(s5, Inches(3.8), Inches(1.7), Inches(2.7), Inches(5.1),
             "Backend API", [
                 "Node.js (Runtime)",
                 "Express.js (REST Framework)",
                 "ESM Modules Structure",
                 "Multer (File Uploads)",
                 "Security Response Headers"
             ])
             
    add_card(s5, Inches(6.8), Inches(1.7), Inches(2.7), Inches(5.1),
             "Database & Cloud", [
                 "MongoDB (NoSQL Document)",
                 "Mongoose ORM (Models)",
                 "Cloudinary (CDN Media)",
                 "Multi-Doc Transactions",
                 "Atomic $inc Operators"
             ])
             
    add_card(s5, Inches(9.8), Inches(1.7), Inches(2.7), Inches(5.1),
             "Security & Payments", [
                 "JWT (Stateless Auth)",
                 "Bcrypt.js (Password Hash)",
                 "Single-Use Reset Tokens",
                 "Razorpay Gateway API",
                 "RBAC Role Protection"
             ])

    # ----------------------------------------------------
    # SLIDE 6: SYSTEM ARCHITECTURE
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "System Architecture & Tiered Structure")
    
    add_card(s6, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Presentation Tier (Client)", [
                 "Customer Portal: Responsive dashboard for purchases, khata ledger balance, and payments.",
                 "Admin POS Portal: Point-of-sale billing, returns, customer profiles, activity feed, and inventory.",
                 "Universal Responsive Engine: Fully tailored for Mobile, Tablet, and Desktop displays."
             ])
             
    add_card(s6, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Application Tier (Server)", [
                 "RESTful API Routes: Organized by domain (/users, /products, /sales, /credits, /payments, /returns).",
                 "Auth & Role Middleware: Protects endpoints via JWT and enforces customer/admin permissions.",
                 "Transaction Manager: Coordinates multi-document ACID write operations."
             ])
             
    add_card(s6, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Data & External Services", [
                 "MongoDB Database: Persists 8 normalized collections (Users, Customers, Products, Sales, etc.).",
                 "Cloudinary Media Cloud: Optimizes image delivery for inventory assets and payment proofs.",
                 "Razorpay Gateway: Secure webhook and signature verification for online payments."
             ])

    # ----------------------------------------------------
    # SLIDE 7: CORE MODULE - POINT OF SALE (POS) BILLING
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Core Module: POS Billing Engine")
    
    add_card(s7, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "POS Billing Workflow", [
                 "Live Catalog Search: Fast keyword lookup with live availability indicators.",
                 "Mixed Cart Items: Supports both cataloged products and unrecorded custom walk-in items.",
                 "Multi-Payment Settlement Modes:",
                 "  • 100% Cash: Instant clearance and cash drawer record.",
                 "  • 100% UPI: Fast digital QR code collection.",
                 "  • 100% Credit (Khata): Automatically binds to customer ledger.",
                 "  • Partial Split: E.g., Pay ₹300 in cash, add ₹700 to credit ledger.",
                 "Printable & Downloadable Receipts with itemized tax and totals."
             ])
             
    add_card(s7, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Real-Time Stock & Atomicity", [
                 "Atomic Stock Deduction: Decrements stock via {$gte: qty} to strictly eliminate overselling.",
                 "Auto Availability Toggle: Marks products 'Out of Stock' when inventory reaches 0.",
                 "Credit Ceiling Guard: Automatically blocks credit sales if customer's debt exceeds Max Borrow Limit.",
                 "Single Transaction Scope: Stock deduction, sale creation, and credit ledger update succeed or fail together."
             ])

    # ----------------------------------------------------
    # SLIDE 8: CORE MODULE - DIGITAL KHATA / CREDIT LEDGER
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Core Module: Digital Credit Ledger ('Khata')")
    
    add_card(s8, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Ledger Mechanics & Tracking", [
                 "Individual Customer Ledger: Tracks total borrowed, amount paid, and net pending balance.",
                 "Due Date Timeline: Configurable repayment deadlines with extension tracking.",
                 "Automatic Ledger Statuses: 'Active', 'Partially Paid', 'Paid', and 'Overdue'.",
                 "Historical Purchase Audit: Full breakdown of which specific sale contributed to the debt."
             ])
             
    add_card(s8, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Trust Score & Risk Control", [
                 "Dynamic Trust Score: Computed based on repayment punctuality and transaction history.",
                 "Manual Borrow Limit Cap: Shop owner can set customized debt ceilings for risky customers.",
                 "Debt Reduction Algorithm: Payments reduce oldest active credit records in FIFO sequence.",
                 "Customer Ledger Visibility: Customers see exact real-time balance on their own devices."
             ])

    # ----------------------------------------------------
    # SLIDE 9: CORE MODULE - PRODUCT RETURNS & REFUNDS
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Core Module: Product Returns & Balance Adjustment")
    
    add_card(s9, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Return Processing Modes", [
                 "Full Sale Return: Returns all items from an invoice, restocks all quantities, and clears associated credit.",
                 "Partial Item Return: Clerk selects specific returned item and quantity (e.g. 2 out of 5 kg returned).",
                 "Return Reasons Tracking: Captures reason (Damaged, Expired, Wrong Item, Customer Change).",
                 "Admin Attribution: Records which store clerk authorized the return."
             ])
             
    add_card(s9, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Automated Financial Balance Correction", [
                 "Automatic Stock Restoration: Returned item quantities are restored back into product inventory.",
                 "Ledger Debt Offset: If sale was on Credit, the return amount automatically subtracts from customer's debt.",
                 "Cash Refund Option: For fully paid sales, calculates direct refund payout to the customer.",
                 "Atomic Execution: Inventory restock and debt reduction occur in an ACID transaction."
             ])

    # ----------------------------------------------------
    # SLIDE 10: CORE MODULE - PAYMENTS & CLAIMS VERIFICATION
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Core Module: Payment Processing & Verification")
    
    add_card(s10, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Customer Payment Channels", [
                 "Cash Claim Submission: Customer submits a claim specifying the store clerk who received physical cash.",
                 "UPI Proof Upload: Customer enters UPI / UTR Transaction ID and uploads payment screenshot.",
                 "Razorpay Gateway: Instant automated settlement via UPI, Debit/Credit Card, or Netbanking.",
                 "Real-Time Statuses: Payments display 'Pending Verification', 'Approved', or 'Rejected'."
             ])
             
    add_card(s10, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Admin Verification Dashboard", [
                 "Dedicated Verification Queue: Lists all pending customer claims awaiting admin approval.",
                 "Receipt Preview: Admins inspect uploaded proof images stored on Cloudinary CDN.",
                 "One-Click Approve / Reject: Approving immediately decreases customer debt via atomic $inc.",
                 "Audit Trail: Saves timestamps, verifying admin name, and remarks for total clarity."
             ])

    # ----------------------------------------------------
    # SLIDE 11: DATABASE DESIGN & SCHEMA
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Database Schema & Entity Relationships")
    
    add_card(s11, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "User & Customer Entities", [
                 "User Collection:",
                 "  • name, username, email, phone",
                 "  • password (bcrypt hash)",
                 "  • role: 'customer' | 'admin' | 'superadmin'",
                 "  • resetPasswordToken & Expire",
                 "Customer Collection:",
                 "  • userId (ref: User)",
                 "  • pendingAmount, totalPurchase",
                 "  • trustScore, manualBorrowLimit"
             ])
             
    add_card(s11, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Catalog & Sales Entities", [
                 "Product Collection:",
                 "  • name, category, price, stock, unit",
                 "  • imageUrl (Cloudinary CDN)",
                 "  • available & deleted flags",
                 "Sale Collection:",
                 "  • customerId, adminId",
                 "  • items: [{ productId, qty, price, total }]",
                 "  • paymentType: cash|upi|credit|partial",
                 "  • paidAmount, pendingAmount, creditId"
             ])
             
    add_card(s11, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Ledger & Returns Entities", [
                 "Credit Collection:",
                 "  • customerId, userId, saleId",
                 "  • borrowedAmount, paidAmount, pending",
                 "  • dueDate, extensionCount, status",
                 "Payment Collection:",
                 "  • creditId, amount, method, status",
                 "  • transactionId, paymentProof, verifiedBy",
                 "Return Collection:",
                 "  • saleId, items, refundAmount, processedBy"
             ])

    # ----------------------------------------------------
    # SLIDE 12: SECURITY & TRANSACTION ACIDITY
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Security Hardening & ACID Atomicity")
    
    add_card(s12, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Database Atomicity & Rollback Protection", [
                 "Multi-Document Transactions: Wrapped in session.withTransaction() so failures trigger complete rollback.",
                 "Zero Orphan Records: If customer balance update fails, sale and payment records are never committed.",
                 "Atomic Math Operators: Uses MongoDB $inc for debt increments/decrements to avoid concurrency races.",
                 "Stock Availability Locks: Conditional updates ensure stock never drops below 0 during peak billing."
             ], bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208), title_color=GREEN)
             
    add_card(s12, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Authentication & Network Security", [
                 "Stateless JWT Verification: Cryptographically signed tokens with secure expiration timestamps.",
                 "Single-Use Password Reset Tokens: Reset tokens are invalidated immediately upon first use in DB.",
                 "Security Response Headers: Hardened with nosniff, frameguard (X-Frame-Options), and XSS filters.",
                 "Environment Isolation: All API keys, Cloudinary secrets, and DB URLs strictly shielded via .env."
             ])

    # ----------------------------------------------------
    # SLIDE 13: ADVANTAGES & BUSINESS IMPACT
    # ----------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Key Advantages & Business Impact")
    
    add_card(s13, Inches(0.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "For Shop Owners (Admins)", [
                 "90% reduction in billing time compared to manual receipt writing.",
                 "Zero revenue leakage from forgotten or disputed customer credit.",
                 "Automated stock alerts prevent product stockouts.",
                 "Audit trail prevents internal staff mishandling."
             ])
             
    add_card(s13, Inches(4.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "For Retail Customers", [
                 "24/7 access to personal digital khata ledger and bills.",
                 "Convenient online repayment options (UPI / QR / Razorpay).",
                 "Total transparency on dues, payments, and discounts.",
                 "No surprise charges or calculation disagreements."
             ])
             
    add_card(s13, Inches(8.8), Inches(1.7), Inches(3.7), Inches(5.1),
             "Technical Superiority", [
                 "100% web-based; works on any browser without installation.",
                 "Zero proprietary hardware requirements.",
                 "Ultra-low infrastructure cost with cloud MongoDB & Cloudinary.",
                 "High scalability for multiple retail branches."
             ])

    # ----------------------------------------------------
    # SLIDE 14: FUTURE ENHANCEMENTS & CONCLUSION
    # ----------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Future Roadmap & Project Conclusion")
    
    add_card(s14, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.1),
             "Future Scope & Enhancements", [
                 "Automated WhatsApp / SMS Reminders for upcoming due dates.",
                 "Thermal POS Receipt Printer & USB Barcode Scanner support.",
                 "Offline Progressive Web App (PWA) billing with IndexedDB sync.",
                 "AI-driven Sales & Inventory demand forecasting.",
                 "Multi-branch Store Management from a single superadmin console."
             ])
             
    add_card(s14, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.1),
             "Conclusion", [
                 "SmartShop successfully bridges the gap between traditional retail credit practices and modern cloud technology.",
                 "Provides a robust, reliable, and user-friendly digital ecosystem for shop owners and customers alike.",
                 "Delivers production-ready stability with verified zero-error frontend and admin builds.",
                 "Demonstrates practical application of fullstack MERN engineering, transaction management, and responsive UI design."
             ])

    # ----------------------------------------------------
    # SLIDE 15: THANK YOU / Q&A (Dark Premium Theme)
    # ----------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    bg15 = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg15.fill.solid()
    bg15.fill.fore_color.rgb = NAVY
    bg15.line.fill.background()

    tbox15 = s15.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf15 = tbox15.text_frame
    tf15.word_wrap = True
    
    p_thx = tf15.paragraphs[0]
    p_thx.text = "Thank You!"
    p_thx.font.name = "Arial"
    p_thx.font.size = Pt(48)
    p_thx.font.bold = True
    p_thx.font.color.rgb = WHITE
    p_thx.alignment = PP_ALIGN.CENTER
    
    p_qa = tf15.add_paragraph()
    p_qa.text = "Questions & Answers (Q&A)"
    p_qa.font.name = "Arial"
    p_qa.font.size = Pt(22)
    p_qa.font.color.rgb = ACCENT
    p_qa.alignment = PP_ALIGN.CENTER
    p_qa.space_before = Pt(12)
    
    p_foot = tf15.add_paragraph()
    p_foot.text = "SmartShop — Digital Shop Management, POS & Khata Ledger Platform"
    p_foot.font.name = "Arial"
    p_foot.font.size = Pt(14)
    p_foot.font.color.rgb = RGBColor(148, 163, 184)
    p_foot.alignment = PP_ALIGN.CENTER
    p_foot.space_before = Pt(24)

    prs.save(pptx_path)
    print("Project Presentation PPTX generated at:", pptx_path)

if __name__ == "__main__":
    create_presentation()
